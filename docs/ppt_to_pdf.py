"""Render the generated deck to PDF.

Run:
    python docs/ppt_to_pdf.py                    # deck -> PDF beside it
    python docs/ppt_to_pdf.py --pptx X --out Y
    python docs/ppt_to_pdf.py --png              # also write one PNG per slide

There is no PowerPoint or LibreOffice on this machine, so the usual
"export as PDF" is not available and the deck cannot even be looked at.
This walks the .pptx with python-pptx and redraws every shape into a PDF
with PyMuPDF.

Two things fall out of doing it this way. The PDF is a real deliverable -
it opens anywhere, which the .pptx does not on a machine without Office.
And because the text is wrapped here with the actual font metrics rather
than estimated, a caption that does not fit is *reported* rather than
guessed at: `--png` then renders each page to an image that can be
inspected directly.

Shape geometry is approximated from the autoshape types the deck actually
uses. It is not a general PowerPoint renderer and does not try to be.
"""

from __future__ import annotations

import argparse
import io
import sys
from pathlib import Path

import pymupdf
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

EMU_PT = 12700.0          # 1 pt = 12700 EMU
LINE = 1.22               # line box, in multiples of the point size

# Arial and Times map onto the base-14 faces MuPDF always has.
FACES = {
    (False, False): "helv", (True, False): "hebo",
    (False, True): "tiro", (True, True): "tibo",
}
_FONTS: dict[str, pymupdf.Font] = {}


def face(bold: bool, serif: bool) -> str:
    return FACES[(bool(bold), bool(serif))]


def font(name: str) -> pymupdf.Font:
    if name not in _FONTS:
        _FONTS[name] = pymupdf.Font(name)
    return _FONTS[name]


def rgb(color) -> tuple[float, float, float] | None:
    """python-pptx colour to a MuPDF triple, or None when it is not solid."""
    try:
        c = color.rgb
    except (AttributeError, TypeError):
        return None
    return (c[0] / 255.0, c[1] / 255.0, c[2] / 255.0)


def fill_of(shape):
    try:
        if shape.fill.type is None or shape.fill.type != 1:   # 1 = solid
            return None
        return rgb(shape.fill.fore_color)
    except (AttributeError, TypeError, ValueError):
        return None


def line_of(shape):
    try:
        col = rgb(shape.line.color)
        if col is None:
            return None, 0.0
        w = shape.line.width
        return col, (w / EMU_PT if w else 1.0)
    except (AttributeError, TypeError, ValueError):
        return None, 0.0


def rotate(points, cx, cy, deg):
    if not deg:
        return points
    import math
    r = math.radians(deg)
    cos, sin = math.cos(r), math.sin(r)
    return [(cx + (x - cx) * cos - (y - cy) * sin,
             cy + (x - cx) * sin + (y - cy) * cos) for x, y in points]


# --------------------------------------------------------------------------
# Geometry
# --------------------------------------------------------------------------

def outline(kind, x, y, w, h):
    """Points for the autoshape types this deck uses. None means 'a rectangle'."""
    # PowerPoint's presets take their inset from the SHORTER side: homePlate
    # and foldedCorner default to 0.16667 of it, rightArrow's head to 0.5.
    # Guessing instead put the pentagon's point three times too deep.
    if kind == MSO_SHAPE.PENTAGON:
        a = min(w, h) * 0.16667
        return [(x, y), (x + w - a, y), (x + w, y + h / 2),
                (x + w - a, y + h), (x, y + h)]
    if kind == MSO_SHAPE.RIGHT_ARROW:
        hw = min(w * 0.5, min(w, h) * 0.5)
        return [(x, y + h * 0.25), (x + w - hw, y + h * 0.25),
                (x + w - hw, y), (x + w, y + h / 2), (x + w - hw, y + h),
                (x + w - hw, y + h * 0.75), (x, y + h * 0.75)]
    if kind == MSO_SHAPE.DOWN_ARROW:
        hh = min(h * 0.5, w)
        return [(x + w * 0.25, y), (x + w * 0.25, y + h - hh),
                (x, y + h - hh), (x + w / 2, y + h), (x + w, y + h - hh),
                (x + w * 0.75, y + h - hh), (x + w * 0.75, y)]
    if kind == MSO_SHAPE.FOLDED_CORNER:
        c = min(w, h) * 0.16667
        return [(x, y), (x + w, y), (x + w, y + h - c), (x + w - c, y + h),
                (x, y + h)]
    return None


def draw_shape(page, shape, x, y, w, h, fill, stroke, width):
    kind = None
    try:
        kind = shape.auto_shape_type
    except (AttributeError, ValueError):
        kind = None
    rot = (shape.rotation or 0) % 360
    cx, cy = x + w / 2, y + h / 2

    if kind == MSO_SHAPE.OVAL:
        page.draw_oval(pymupdf.Rect(x, y, x + w, y + h), color=stroke,
                       fill=fill, width=width)
        return
    if kind == MSO_SHAPE.DONUT:
        page.draw_oval(pymupdf.Rect(x, y, x + w, y + h), color=None, fill=fill)
        # The hole. PowerPoint's default is half the diameter.
        i = min(w, h) * 0.25
        page.draw_oval(pymupdf.Rect(x + i, y + i, x + w - i, y + h - i),
                       color=None, fill=(1, 1, 1))
        return
    if kind == MSO_SHAPE.CAN:
        cap = min(h * 0.22, w * 0.5)
        page.draw_rect(pymupdf.Rect(x, y + cap / 2, x + w, y + h - cap / 2),
                       color=None, fill=fill)
        page.draw_oval(pymupdf.Rect(x, y, x + w, y + cap), color=None, fill=fill)
        page.draw_oval(pymupdf.Rect(x, y + h - cap, x + w, y + h),
                       color=None, fill=fill)
        return

    pts = outline(kind, x, y, w, h)
    if pts:
        pts = rotate(pts, cx, cy, rot)
        page.draw_polyline([pymupdf.Point(*p) for p in pts] +
                           [pymupdf.Point(*pts[0])],
                           color=stroke, fill=fill, width=width, closePath=True)
        return

    radius = 0.16667 if kind == MSO_SHAPE.ROUNDED_RECTANGLE else None
    rect = pymupdf.Rect(x, y, x + w, y + h)
    if radius and min(w, h) > 2:
        page.draw_rect(rect, color=stroke, fill=fill, width=width,
                       radius=min(radius, 0.49))
    else:
        page.draw_rect(rect, color=stroke, fill=fill, width=width)


# --------------------------------------------------------------------------
# Text
# --------------------------------------------------------------------------

# The base-14 faces cannot encode these, and PyMuPDF silently substitutes a
# dot. The deck keeps the proper glyphs for PowerPoint; only the PDF flattens.
SUBS = str.maketrans({
    "“": '"', "”": '"', "‘": "'", "’": "'",
    "–": "-", "—": "-", "→": "->", "←": "<-",
    "σ": "sd", "◆": "*", "•": "-", "·": "-",
})


def wrap(text, fnt, size, width):
    text = text.translate(SUBS)
    words = text.split()
    if not words:
        return [""]
    lines, cur = [], ""
    for word in words:
        trial = f"{cur} {word}".strip()
        if not cur or fnt.text_length(trial, fontsize=size) <= width:
            cur = trial
        else:
            lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    return lines


def paragraphs_of(shape, default_size, default_color, serif):
    out = []
    for para in shape.text_frame.paragraphs:
        text = "".join(r.text for r in para.runs) or para.text
        if not text.strip():
            continue
        run = para.runs[0] if para.runs else None
        size = (run.font.size.pt if run is not None and run.font.size
                else default_size)
        bold = bool(run.font.bold) if run is not None and run.font.bold else False
        col = rgb(run.font.color) if run is not None else None
        before = para.space_before.pt if para.space_before else 0.0
        out.append({
            "text": text, "size": size, "bold": bold,
            "color": col or default_color, "align": para.alignment,
            "before": before, "serif": serif,
        })
    return out


def draw_text(page, shape, x, y, w, h, *, default_size=11.0,
              default_color=(0.2, 0.25, 0.31), serif=False, report=None):
    tf = shape.text_frame
    paras = paragraphs_of(shape, default_size, default_color, serif)
    if not paras:
        return

    pad_l = (tf.margin_left / EMU_PT) if tf.margin_left is not None else 7.2
    pad_r = (tf.margin_right / EMU_PT) if tf.margin_right is not None else 7.2
    pad_t = (tf.margin_top / EMU_PT) if tf.margin_top is not None else 3.6
    avail = max(4.0, w - pad_l - pad_r)

    for p in paras:
        p["lines"] = wrap(p["text"], font(face(p["bold"], p["serif"])),
                          p["size"], avail)

    total = sum(p["before"] + len(p["lines"]) * p["size"] * LINE for p in paras)
    if report is not None and total > h - pad_t + 0.5:
        report.append((round(total - (h - pad_t), 1), shape.name,
                       paras[0]["text"][:44]))

    anchor = tf.vertical_anchor
    if anchor == MSO_ANCHOR.MIDDLE:
        cursor = y + (h - total) / 2
    elif anchor == MSO_ANCHOR.BOTTOM:
        cursor = y + h - total - pad_t
    else:
        cursor = y + pad_t

    for p in paras:
        cursor += p["before"]
        fnt_name = face(p["bold"], p["serif"])
        fnt = font(fnt_name)
        for line in p["lines"]:
            lw = fnt.text_length(line, fontsize=p["size"])
            if p["align"] == PP_ALIGN.LEFT:
                lx = x + pad_l
            elif p["align"] == PP_ALIGN.RIGHT:
                lx = x + w - pad_r - lw
            else:
                lx = x + (w - lw) / 2
            baseline = cursor + p["size"] * 0.98
            page.insert_text(pymupdf.Point(lx, baseline), line,
                             fontname=fnt_name, fontsize=p["size"],
                             color=p["color"])
            cursor += p["size"] * LINE


# --------------------------------------------------------------------------

def render_table(page, shape, x, y, w, h):
    table = shape.table
    widths = [c.width / EMU_PT for c in table.columns]
    heights = [r.height / EMU_PT for r in table.rows]
    scale_w = w / sum(widths) if sum(widths) else 1
    scale_h = h / sum(heights) if sum(heights) else 1
    cy = y
    for ri, row in enumerate(table.rows):
        cx = x
        rh = heights[ri] * scale_h
        for ci, cell in enumerate(row.cells):
            cw = widths[ci] * scale_w
            fill = fill_of(cell)
            page.draw_rect(pymupdf.Rect(cx, cy, cx + cw, cy + rh),
                           color=(0.83, 0.87, 0.92), fill=fill, width=0.5)
            draw_text(page, cell, cx, cy, cw, rh, default_size=9.0)
            cx += cw
        cy += rh


def render(pptx_path: Path, out: Path, png: bool) -> None:
    prs = Presentation(str(pptx_path))
    pw = prs.slide_width / EMU_PT
    ph = prs.slide_height / EMU_PT
    doc = pymupdf.open()
    overflow = []

    for index, slide in enumerate(prs.slides, 1):
        page = doc.new_page(width=pw, height=ph)
        page.draw_rect(pymupdf.Rect(0, 0, pw, ph), color=None, fill=(1, 1, 1))
        for shape in slide.shapes:
            if shape.left is None or shape.width is None:
                continue
            x, y = shape.left / EMU_PT, shape.top / EMU_PT
            w, h = shape.width / EMU_PT, shape.height / EMU_PT
            kind = str(shape.shape_type)

            if kind.startswith("PICTURE"):
                try:
                    page.insert_image(pymupdf.Rect(x, y, x + w, y + h),
                                      stream=io.BytesIO(shape.image.blob))
                except (AttributeError, ValueError, RuntimeError) as exc:
                    print(f"  slide {index}: image skipped ({exc})",
                          file=sys.stderr)
                continue

            if kind.startswith("TABLE"):
                render_table(page, shape, x, y, w, h)
                continue

            if kind.startswith("LINE"):
                col, lw = line_of(shape)
                page.draw_line(pymupdf.Point(x, y),
                               pymupdf.Point(x + w, y + h),
                               color=col or (0.83, 0.87, 0.92),
                               width=lw or 1.0)
                continue

            has_text = shape.has_text_frame and shape.text_frame.text.strip()
            fill = fill_of(shape)
            stroke, lw = line_of(shape)
            if fill is not None or stroke is not None:
                draw_shape(page, shape, x, y, w, h, fill, stroke, lw or 1.0)

            if has_text:
                # Template furniture carries theme colours python-pptx cannot
                # resolve; give the title the deck's navy and the rest grey.
                name = shape.name.lower()
                if "title" in name and "sub" not in name:
                    draw_text(page, shape, x, y, w, h, default_size=28.0,
                              default_color=(0.12, 0.29, 0.49), serif=True,
                              report=overflow)
                elif "footer" in name or "slide number" in name:
                    draw_text(page, shape, x, y, w, h, default_size=10.0,
                              default_color=(0.55, 0.58, 0.62))
                else:
                    draw_text(page, shape, x, y, w, h, report=overflow)

    out.parent.mkdir(parents=True, exist_ok=True)
    # Each draw_* call appends its own content stream and leaves the previous
    # objects behind; saved as-is the deck came to 300 MB of garbage for 2 MB
    # of drawing.
    doc.save(str(out), garbage=4, deflate=True, clean=True)
    print(f"wrote {out}  ({len(doc)} pages, {out.stat().st_size // 1024} KB)")

    if png:
        for i, page in enumerate(doc, 1):
            pix = page.get_pixmap(dpi=110)
            shot = out.with_name(f"{out.stem}-{i}.png")
            pix.save(str(shot))
            print(f"  {shot.name}  {pix.width}x{pix.height}")

    if overflow:
        print("\ntext that did not fit its shape (overflow pt, shape, text):",
              file=sys.stderr)
        for row in sorted(overflow, reverse=True):
            print("   ", row, file=sys.stderr)
    else:
        print("all text fitted its shape")
    doc.close()


def main() -> None:
    root = Path(__file__).parent.parent
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--pptx", type=Path,
                    default=root / "SIH2026_PS26153_Presentation.pptx")
    ap.add_argument("--out", type=Path, default=None)
    ap.add_argument("--png", action="store_true",
                    help="also write one PNG per slide, to look at")
    args = ap.parse_args()

    if not args.pptx.exists():
        raise SystemExit(f"Deck not found: {args.pptx}")
    render(args.pptx, args.out or args.pptx.with_suffix(".pdf"), args.png)


if __name__ == "__main__":
    main()
