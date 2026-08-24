"""Fill the official SIH 2026 idea-submission template with our content.

Run:
    python docs/make_ppt.py                     # uses the template in Downloads
    python docs/make_ppt.py --template X --out Y

The template's own rules drive every layout decision here:

  - **Six slides maximum, title slide included.** So the instructions slide is
    deleted and nothing new is added: exactly the five content sections the
    template defines, and no appendix.
  - **"Avoid paragraphs, use points / diagrams / infographics."** Content is
    laid out as labelled cards, a pipeline diagram and a results table. There
    is no running prose anywhere on the deck.
  - **"Only use the provided template without changing the idea details
    pointers."** Section titles are untouched. Each section's guidance pointers
    become the card headings, so a reviewer can still see the required
    structure - we answer the pointers rather than replacing them.

Colours and fonts are read off the template's own theme (Times New Roman
titles, Arial body, #1F497D / #4F81BD) so added shapes look native rather than
pasted in.
"""

from __future__ import annotations

import argparse
import math
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --------------------------------------------------------------------------
# Design tokens, taken from the template theme
# --------------------------------------------------------------------------

NAVY = RGBColor(0x1F, 0x49, 0x7D)      # theme dk1 - headings, dark bands
BLUE = RGBColor(0x4F, 0x81, 0xBD)      # theme accent1 - primary cards
RED = RGBColor(0xC0, 0x50, 0x4D)       # theme accent2 - risks, the problem
GREEN = RGBColor(0x77, 0x93, 0x3C)     # theme accent3, darkened for contrast
SLATE = RGBColor(0x33, 0x3F, 0x50)     # body text
MUTED = RGBColor(0x5A, 0x64, 0x78)     # secondary text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xF2, 0xF5, 0xFA)   # very light blue-grey card fill
CARD_LINE = RGBColor(0xD3, 0xDD, 0xEA)

BODY_FONT = "Arial"

# Usable content band, in inches. Above 1.35 sits the title and the team oval;
# below 6.85 sits the template's footer bar.
TOP = 1.42
BOTTOM = 6.80
LEFT = 0.45
RIGHT = 12.85
WIDTH = RIGHT - LEFT


# --------------------------------------------------------------------------
# Shape helpers
# --------------------------------------------------------------------------


def clear(shape) -> None:
    """Empty a text frame without destroying the shape."""
    tf = shape.text_frame
    tf.clear()
    tf.paragraphs[0].text = ""


def set_text(tf, lines, *, size=12, color=SLATE, bold=False, space_after=4,
             align=PP_ALIGN.LEFT, bullet_char="• "):
    """Write lines into a text frame.

    Args:
        lines: list of str, or (text, level) tuples. Level 0 gets a bullet,
            level 1 is an unbulleted sub-line, level 2 is a bold sub-heading.
    """
    tf.clear()
    tf.word_wrap = True
    first = True

    for item in lines:
        text, level = item if isinstance(item, tuple) else (item, 0)
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = align
        para.space_after = Pt(space_after)

        run = para.add_run()
        if level == 0:
            run.text = f"{bullet_char}{text}"
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.bold = bold
        elif level == 1:
            run.text = f"    {text}"
            run.font.size = Pt(size - 1)
            run.font.color.rgb = MUTED
        else:
            run.text = text
            run.font.size = Pt(size + 1)
            run.font.color.rgb = NAVY
            run.font.bold = True
        run.font.name = BODY_FONT


def add_card(slide, x, y, w, h, title, lines, *, accent=BLUE, size=11,
             title_size=12.5):
    """A titled card: coloured header strip over a light body panel."""
    head_h = 0.36
    head = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(head_h))
    head.fill.solid()
    head.fill.fore_color.rgb = accent
    head.line.fill.background()
    head.shadow.inherit = False
    tf = head.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.10)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = BODY_FONT

    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y + head_h),
                                  Inches(w), Inches(h - head_h))
    body.fill.solid()
    body.fill.fore_color.rgb = CARD_BG
    body.line.color.rgb = CARD_LINE
    body.line.width = Pt(0.75)
    body.shadow.inherit = False
    btf = body.text_frame
    # Autoshapes default to middle anchoring, which parks a short list in the
    # vertical centre of its card and leaves a dead band above it.
    btf.vertical_anchor = MSO_ANCHOR.TOP
    btf.margin_left = Inches(0.13)
    btf.margin_right = Inches(0.10)
    btf.margin_top = Inches(0.11)
    set_text(btf, lines, size=size)
    return body


def set_title(slide, text: str) -> None:
    """Replace a title's text while keeping the template's title formatting.

    The title frame holds several runs, so overwriting only the first leaves
    the rest of the placeholder text visible underneath. Capture the font off
    run 0, clear the frame, then rebuild a single run with it.
    """
    tf = find(slide, "Title 1").text_frame
    src = tf.paragraphs[0].runs[0].font
    name, size, bold = src.name, src.size, src.bold
    tf.clear()
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.name = name
    run.font.size = size
    run.font.bold = bold


def add_band(slide, x, y, w, h, text, *, fill=NAVY, size=14, color=WHITE,
             align=PP_ALIGN.CENTER, bold=True):
    """A full-width statement band."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = BODY_FONT
    return box


def add_step(slide, x, y, w, h, headline, detail, *, fill=BLUE):
    """One box in the pipeline diagram."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = headline
    r.font.size = Pt(10.5)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = BODY_FONT

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = detail
    r2.font.size = Pt(8.5)
    r2.font.color.rgb = RGBColor(0xE4, 0xEC, 0xF7)
    r2.font.name = BODY_FONT
    return box


def add_arrow(slide, x, y, w, h):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = RGBColor(0xB8, 0xC6, 0xDA)
    arr.line.fill.background()
    arr.shadow.inherit = False
    return arr


def add_table(slide, x, y, w, h, rows_data, col_widths=None, *, header_fill=NAVY,
              size=10, highlight_row=None):
    n_rows, n_cols = len(rows_data), len(rows_data[0])
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    table = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)

    for ri, row in enumerate(rows_data):
        table.rows[ri].height = Inches(h / n_rows)
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.margin_left = Inches(0.07)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = header_fill
            elif ri == highlight_row:
                cell.fill.fore_color.rgb = RGBColor(0xE3, 0xEE, 0xE0)
            else:
                cell.fill.fore_color.rgb = WHITE

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(size)
            r.font.name = BODY_FONT
            r.font.bold = (ri == 0) or (ri == highlight_row)
            r.font.color.rgb = WHITE if ri == 0 else SLATE
    return table


def delete_slide(prs, index: int) -> None:
    """python-pptx has no public slide delete; drop the id and the relationship."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rid = slides[index].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    prs.part.drop_rel(rid)
    xml_slides.remove(slides[index])


def find(slide, name_startswith):
    for sh in slide.shapes:
        if sh.name.startswith(name_startswith):
            return sh
    return None


# --------------------------------------------------------------------------
# Diagram primitives
# --------------------------------------------------------------------------
# The template asks for "points / diagrams / infographics", not paragraphs, and
# a deck read from the back of a room is carried by shape and colour long
# before anybody parses a sentence. These build the recurring figures - mind
# map, arrow flow, chain, shields, hangers, pillars, hub, ring - so each slide
# is a picture with labels rather than three columns of prose.

ORANGE = RGBColor(0xE3, 0x7B, 0x2A)
PURPLE = RGBColor(0x7A, 0x45, 0xD1)
TEAL = RGBColor(0x11, 0x83, 0x8B)
PINK = RGBColor(0xC2, 0x3B, 0x77)
AMBER = RGBColor(0xB8, 0x7D, 0x00)
PALETTE = [BLUE, TEAL, ORANGE, PURPLE, GREEN, PINK]


def _flat(shape, fill, line=None, line_w=1.0):
    """Flat fill, no shadow. PowerPoint's default shadow reads as clip art."""
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    return shape


def _write(tf, lines, *, size=10, color=WHITE, bold=False, align=PP_ALIGN.CENTER,
           anchor=MSO_ANCHOR.MIDDLE, spacing=0):
    """Write into a text frame. Each line is a str or (text, size, bold)."""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    norm = [(l if isinstance(l, tuple) else (l, size, bold)) for l in lines]
    for i, (text, sz, bd) in enumerate(norm):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        if spacing and i:
            para.space_before = Pt(spacing)
        run = para.add_run()
        run.text = text
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.color.rgb = color
        run.font.name = BODY_FONT
    return tf


def shp(slide, kind, x, y, w, h, *, fill=BLUE, line=None, line_w=1.0, lines=(),
        size=10, color=WHITE, bold=False, anchor=MSO_ANCHOR.MIDDLE,
        align=PP_ALIGN.CENTER):
    sh = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    _flat(sh, fill, line, line_w)
    if lines:
        _write(sh.text_frame, lines, size=size, color=color, bold=bold,
               align=align, anchor=anchor)
    else:
        sh.text_frame.text = ""
    return sh


def label(slide, x, y, w, h, lines, *, size=9, color=SLATE, bold=False,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, spacing=0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _write(tb.text_frame, lines, size=size, color=color, bold=bold,
           align=align, anchor=anchor, spacing=spacing)
    return tb


def connect(slide, x1, y1, x2, y2, *, color=CARD_LINE, width=1.0, dash=True):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),
                                    Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    if dash:
        ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return ln


def section(slide, x, y, w, text, *, color=NAVY, size=12.5):
    """A heading with a colour rule under it, as on the reference deck."""
    label(slide, x, y, w, 0.28, [text], size=size, color=color, bold=True,
          align=PP_ALIGN.LEFT)
    return shp(slide, MSO_SHAPE.RECTANGLE, x + 0.04, y + 0.28, w - 0.08, 0.035,
               fill=color)


def down_pentagon(slide, x, y, w, h, *, fill=BLUE):
    """A pentagon pointing down, fitted to the box (x, y, w, h).

    PowerPoint's pentagon points right, so it is built on its side and rotated
    about its own centre - which means the pre-rotation box has w and h
    swapped, or the shape lands outside the slot it was given.
    """
    cx, cy = x + w / 2, y + h / 2
    sh = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(cx - h / 2),
                                Inches(cy - w / 2), Inches(h), Inches(w))
    _flat(sh, fill)
    sh.rotation = 90
    sh.text_frame.text = ""
    return sh


def add_mindmap(slide, x, y, w, h, centre, left_items, right_items):
    """Centre node with labelled branches either side, joined by dashed leads."""
    col_w = w * 0.30
    mid_w = w - 2 * col_w - 0.44
    mid_x = x + col_w + 0.22
    cx, cy = mid_x + mid_w / 2, y + h / 2

    for items, bx, is_left in ((left_items, x, True),
                               (right_items, x + w - col_w, False)):
        n = len(items)
        slot = h / n
        anchor_x = mid_x if is_left else mid_x + mid_w
        for i, (head, subs, colour) in enumerate(items):
            by = y + i * slot + (slot - 0.94) / 2
            connect(slide, anchor_x, cy, bx + (col_w if is_left else 0), by + 0.22)
            shp(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, col_w, 0.42,
                fill=colour, lines=[head], size=9.5, bold=True)
            label(slide, bx, by + 0.46, col_w, 0.48,
                  [("- " + s, 7.5, False) for s in subs], color=MUTED)

    shp(slide, MSO_SHAPE.ROUNDED_RECTANGLE, mid_x, cy - 0.46, mid_w, 0.92,
        fill=NAVY, lines=[(centre[0], 12, True), (centre[1], 8, False)],
        color=WHITE)


def add_arrow_flow(slide, x, y, w, items):
    """A wide arrow carrying numbered stops, labels alternating above/below."""
    shaft_h = 0.44
    ay = y + 0.82
    shp(slide, MSO_SHAPE.RIGHT_ARROW, x, ay, w, shaft_h, fill=CARD_BG)
    n = len(items)
    inner = w - 0.52
    box_w = min(0.86, inner / n - 0.12)
    step = (inner - box_w) / max(1, n - 1)
    for i, (head, detail) in enumerate(items):
        bx = x + 0.10 + i * step
        colour = PALETTE[i % len(PALETTE)]
        shp(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, ay - 0.09, box_w,
            shaft_h + 0.18, fill=colour, lines=[str(i + 1)], size=13, bold=True)
        above = i % 2 == 0
        cap_w = box_w + 1.04
        lx = min(max(bx - 0.52, x - 0.04), x + w - cap_w + 0.04)
        label(slide, lx, (y - 0.06) if above else (ay + 0.50),
              cap_w, 0.64,
              [(head, 8.5, True), (detail, 7, False)],
              color=colour if above else SLATE, spacing=1)


def add_chain(slide, x, y, w, h, items):
    """Interlocking outlines - the reference deck's uniqueness figure."""
    n = len(items)
    link_w = (w + (n - 1) * 0.24) / n
    step = link_w - 0.24
    for i, (head, detail) in enumerate(items):
        colour = PALETTE[(i + 2) % len(PALETTE)]
        lx = x + i * step
        shp(slide, MSO_SHAPE.ROUNDED_RECTANGLE, lx, y, link_w, h,
            fill=None, line=colour, line_w=4.5)
        # Inset past the interlock: the links overlap by design, so captions
        # set to the full link width would sit on each other.
        label(slide, lx + 0.20, y + h + 0.05, link_w - 0.40, 0.72,
              [(head, 8, True), (detail, 7, False)], color=colour, spacing=1)


def add_shields(slide, x, y, w, items, *, h=1.24):
    """Pentagon shields in a row - the feasibility dimensions."""
    n = len(items)
    gap = 0.11
    sw = (w - (n - 1) * gap) / n
    for i, (head, detail) in enumerate(items):
        colour = PALETTE[i % len(PALETTE)]
        sx = x + i * (sw + gap)
        down_pentagon(slide, sx, y, sw, h, fill=colour)
        label(slide, sx, y + 0.14, sw, 0.86,
              [(head, 9, True), (detail, 7, False)], color=WHITE, spacing=2)


def add_hangers(slide, x, y, w, items):
    """Boxes suspended from a rail - the reference deck's risk figure."""
    n = len(items)
    gap = 0.10
    bw = (w - (n - 1) * gap) / n
    for i, (head, detail) in enumerate(items):
        colour = PALETTE[i % len(PALETTE)]
        bx = x + i * (bw + gap)
        connect(slide, bx + bw / 2, y, bx + bw / 2, y + 0.24, color=MUTED,
                width=1.25, dash=False)
        shp(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, y + 0.24, bw, 0.44,
            fill=colour, lines=[head], size=8.5, bold=True)
        label(slide, bx, y + 0.72, bw, 0.66, [(detail, 7, False)], color=MUTED)


def add_pillars(slide, x, y, w, items):
    """Cylinders - the mitigations holding those risks up."""
    n = len(items)
    gap = 0.11
    pw = (w - (n - 1) * gap) / n
    for i, (head, detail) in enumerate(items):
        colour = PALETTE[i % len(PALETTE)]
        px = x + i * (pw + gap)
        shp(slide, MSO_SHAPE.CAN, px + pw * 0.14, y, pw * 0.72, 0.56,
            fill=colour)
        label(slide, px, y + 0.60, pw, 0.70,
              [(head, 8.5, True), (detail, 7, False)], color=colour, spacing=1)


def add_hub(slide, cx, cy, centre, items, *, rx=2.55, ry=1.62):
    """Hub and spokes - who the work lands on."""
    n = len(items)
    for i, (head, detail) in enumerate(items):
        ang = -math.pi / 2 + i * 2 * math.pi / n
        bx = cx + rx * math.cos(ang) - 0.88
        by = cy + ry * math.sin(ang) - 0.33
        colour = PALETTE[i % len(PALETTE)]
        connect(slide, cx + 0.62 * math.cos(ang), cy + 0.62 * math.sin(ang),
                bx + 0.88, by + 0.33, color=colour, width=1.25, dash=False)
        shp(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, 1.76, 0.66,
            fill=colour, lines=[(head, 9, True), (detail, 7, False)], color=WHITE)
    shp(slide, MSO_SHAPE.OVAL, cx - 0.70, cy - 0.70, 1.40, 1.40, fill=NAVY,
        lines=[(centre[0], 10.5, True), (centre[1], 7.5, False)], color=WHITE)


def add_ring(slide, cx, cy, items, *, r=0.92):
    """A ring with claims set around it.

    Built from a donut plus a marker per claim rather than from pie wedges:
    wedge geometry is driven by shape adjustments whose units differ between
    renderers, and a deck that only lays out correctly in one viewer is not
    worth the extra colour.
    """
    shp(slide, MSO_SHAPE.DONUT, cx - r, cy - r, 2 * r, 2 * r, fill=CARD_BG)
    n = len(items)
    for i, (head, detail) in enumerate(items):
        ang = -math.pi / 2 + i * 2 * math.pi / n
        colour = PALETTE[i % len(PALETTE)]
        mx = cx + r * 0.86 * math.cos(ang)
        my = cy + r * 0.86 * math.sin(ang)
        shp(slide, MSO_SHAPE.OVAL, mx - 0.15, my - 0.15, 0.30, 0.30,
            fill=colour, lines=[str(i + 1)], size=8, bold=True)
        lx = cx + (r + 1.34) * math.cos(ang) - 0.94
        ly = cy + (r + 0.78) * math.sin(ang) - 0.30
        label(slide, lx, ly, 1.88, 0.62,
              [(head, 8.5, True), (detail, 7, False)], color=colour, spacing=1)


def add_stack_table(slide, x, y, w, rows, *, row_h=0.40, gap=0.055):
    """The reference deck's tech-stack grid: coloured label, plain value."""
    lab_w = w * 0.44
    for i, (head, value) in enumerate(rows):
        colour = PALETTE[i % len(PALETTE)]
        ry = y + i * (row_h + gap)
        shp(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, ry, lab_w, row_h,
            fill=colour, lines=[head], size=8.5, bold=True, align=PP_ALIGN.LEFT)
        shp(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + lab_w + 0.07, ry,
            w - lab_w - 0.07, row_h, fill=CARD_BG, lines=[value], size=8.5,
            color=SLATE, align=PP_ALIGN.LEFT)


def add_stack_layers(slide, x, y, w, layers, *, band_h=0.50, arrow_h=0.15):
    """A stacked architecture figure: layer bands joined by down arrows."""
    cy = y
    for i, (head, detail, colour) in enumerate(layers):
        shp(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, w, band_h, fill=colour,
            lines=[(head + "   ", 9.5, True), (detail, 7.5, False)], color=WHITE)
        cy += band_h
        if i < len(layers) - 1:
            shp(slide, MSO_SHAPE.DOWN_ARROW, x + w / 2 - 0.10, cy + 0.015,
                0.20, arrow_h - 0.03, fill=CARD_LINE)
            cy += arrow_h
    return cy


def add_notes(slide, x, y, w, items, *, h=0.72):
    """Folded-corner notes - the reference deck's team/reference strip."""
    n = len(items)
    gap = 0.10
    nw = (w - (n - 1) * gap) / n
    for i, text in enumerate(items):
        colour = PALETTE[i % len(PALETTE)]
        shp(slide, MSO_SHAPE.FOLDED_CORNER, x + i * (nw + gap), y, nw, h,
            fill=colour, lines=[text], size=8, bold=True)


# --------------------------------------------------------------------------
# Slide builders
# --------------------------------------------------------------------------


def build_title(slide, team_id: str, team_name: str) -> None:
    box = find(slide, "TextBox 9")
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    rows = [
        ("Problem Statement ID", "26153"),
        ("Problem Statement Title",
         "AI based Network Attack Forecasting from Network Traffic Data"),
        ("Theme", "Blockchain & Cybersecurity"),
        ("PS Category", "Software"),
        ("Organisation", "National Technical Research Organisation (NTRO)"),
        ("Team ID", team_id),
        ("Team Name", team_name),
    ]
    # Drop any row we have no real value for. A line reading "[Team ID]" looks
    # like an unfilled template; an absent line just looks finished, and the ID
    # can be added later with --team-id.
    rows = [(k, v) for k, v in rows if v and not v.startswith("[")]

    for i, (k, v) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        rk = p.add_run()
        rk.text = f"{k} – "
        rk.font.size = Pt(13)
        rk.font.bold = True
        rk.font.color.rgb = NAVY
        rk.font.name = BODY_FONT
        rv = p.add_run()
        rv.text = v
        rv.font.size = Pt(13)
        rv.font.color.rgb = SLATE
        rv.font.name = BODY_FONT


def build_idea(slide) -> None:
    set_title(slide, "PREDICTIVE CYBER DEFENCE")
    clear(find(slide, "TextBox 8"))

    add_band(slide, LEFT, TOP, WIDTH, 0.48,
             "We do not classify traffic - we learn how a network's state "
             "evolves, then run that model forward to see where a host is heading.",
             size=12.5)

    lw = 5.90
    rx = LEFT + lw + 0.55
    rw = RIGHT - rx

    # ---- left: what the system is made of -------------------------------
    section(slide, LEFT, TOP + 0.60, lw, "AttackForecast - what it is made of")
    add_mindmap(
        slide, LEFT, TOP + 1.02, lw, 3.86,
        ("AttackForecast", "RSSM world model  -  138k parameters"),
        [
            ("Flow features", ["39 per 60-second window",
                               "bytes, ports, entropy"], BLUE),
            ("Packet features", ["23 from raw PCAP",
                                 "TTL variance, TCP window"], TEAL),
            ("World model", ["encoder + GRU + prior",
                             "learns state transitions"], ORANGE),
        ],
        [
            ("Forecasting", ["10 states imagined ahead",
                             "uncertainty band"], PURPLE),
            ("ATT&CK mapping", ["5 kill-chain stages",
                                "technique IDs attached"], GREEN),
            ("Explainability", ["integrated gradients",
                                "temporal attention"], PINK),
        ],
    )

    # ---- right: how it answers the problem, and what is new -------------
    section(slide, rx, TOP + 0.60, rw, "How it addresses the problem")
    add_arrow_flow(slide, rx, TOP + 1.10, rw, [
        ("Observe", "60-second host windows"),
        ("Learn transitions", "what state follows this state"),
        ("Cut the traffic off", "model sees nothing more"),
        ("Imagine 10 steps", "it generates the future itself"),
        ("Score each state", "risk + ATT&CK stage"),
    ])

    section(slide, rx, TOP + 3.20, rw, "Innovation and uniqueness")
    add_chain(slide, rx, TOP + 3.63, rw, 0.50, [
        ("World model, not a classifier",
         "heads read the latent state, so they work on imagined "
         "states"),
        ("Falsifiable by construction",
         "imagination-only loss gives zero encoder gradient - "
         "in our tests"),
        ("Runs offline on a CPU",
         "no cloud, no GPU, no external API - air-gap ready for "
         "CII networks"),
    ])

    add_band(slide, LEFT, BOTTOM - 0.42, WIDTH, 0.42,
             "Working prototype today  -  258,229 network states from all 13 "
             "CTU-13 captures  -  28 of 30 infected hosts flagged at 4.9% "
             "false alarms",
             fill=GREEN, size=11)


def build_technical(slide) -> None:
    clear(find(slide, "TextBox 8"))

    lw = 7.05
    rx = LEFT + lw + 0.55
    rw = RIGHT - rx

    section(slide, LEFT, TOP, lw, "Architecture - traffic in, forecast out")
    add_stack_layers(slide, LEFT + 0.30, TOP + 0.50, lw - 0.60, [
        ("INGEST", "NetFlow / .binetflow  -  raw PCAP via Scapy", NAVY),
        ("FEATURE EXTRACTION", "39 flow features  +  23 packet features", BLUE),
        ("NETWORK STATE", "one cell per (host, 60-second window)", TEAL),
        ("ENCODER + GRU", "latent state carrying the host's history", ORANGE),
        ("PRIOR ROLLOUT", "10 states imagined with no traffic observed", PURPLE),
        ("PREDICTION HEADS", "infiltration risk  -  ATT&CK stage  -  surprise", PINK),
        ("DECISION SUPPORT", "dashboard, replay, three explanation channels", GREEN),
    ])

    section(slide, rx, TOP, rw, "Tech stack used")
    add_stack_table(slide, rx, TOP + 0.50, rw, [
        ("Language", "Python 3.13"),
        ("Deep learning", "PyTorch (CPU-only)"),
        ("Packet parsing", "Scapy - raw PCAP, flow rebuild"),
        ("Data pipeline", "pandas + NumPy, vectorised"),
        ("Baselines & metrics", "scikit-learn"),
        ("Backend", "FastAPI + Uvicorn"),
        ("Frontend", "Single-file HTML, zero CDN"),
        ("Explainability", "Integrated Gradients + attention"),
        ("Dataset", "CTU-13 - 13 real captures, CC BY"),
        ("Deployment", "Fully offline, no external API"),
    ])

    add_band(slide, LEFT, BOTTOM - 0.44, WIDTH, 0.44,
             "Validation: split by time, never randomly  -  baselines given "
             "identical features, scaler and labels  -  seeded rollouts, so the "
             "benchmark reproduces byte for byte",
             fill=RED, size=11)


def build_feasibility(slide, images_dir: Path) -> None:
    clear(find(slide, "TextBox 8"))

    lw = 6.00
    rx = LEFT + lw + 0.40
    rw = RIGHT - rx

    # ---- left: is it deployable, and does it actually run ---------------
    section(slide, LEFT, TOP, lw, "Feasibility of AttackForecast")
    add_shields(slide, LEFT, TOP + 0.48, lw, [
        ("Technical", "CPU-only, 138k params"),
        ("Financial", "zero licence cost"),
        ("Operational", "consumes NetFlow already exported"),
        ("Legal", "data never leaves the network"),
        ("Social", "analyst keeps the decision"),
    ])

    section(slide, LEFT, TOP + 1.92, lw, "Live output from the running prototype")
    shot = images_dir / "g-topology.png"
    if shot.exists():
        with Image.open(shot) as img:
            aspect = img.width / img.height
        # Sized from the space left above the footer, then from the crop's
        # own aspect - a squashed screenshot looks exactly like one.
        top = TOP + 2.38
        sh_h = min(BOTTOM - 0.34 - top, 3.95 / aspect)
        sw = sh_h * aspect
        slide.shapes.add_picture(
            str(shot), Inches(LEFT + (lw - sw) / 2), Inches(top),
            width=Inches(sw), height=Inches(sh_h),
        )
        label(slide, LEFT, top + sh_h + 0.04, lw, 0.28,
              ["3D topology - the compromised host and its malicious fan-out"],
              size=8, color=MUTED)
    else:
        print(f"  warning: {shot.name} missing, slide 4 screenshot skipped",
              file=sys.stderr)

    # ---- right: what could go wrong, and what we did about it -----------
    section(slide, rx, TOP, rw, "Potential challenges and risks", color=RED)
    add_hangers(slide, rx, TOP + 0.48, rw, [
        ("Label leak", "a metadata flag became a perfect label proxy"),
        ("Detection is only a tie", "0.979 against 0.977 for logistic regression"),
        ("Short-burst compromise", "the risk head needs sustained activity"),
        ("Stages do not cross families", "progression did not transfer"),
    ])

    section(slide, rx, TOP + 1.96, rw, "Strategies for overcoming these", color=GREEN)
    add_pillars(slide, rx, TOP + 2.44, rw, [
        ("Found and removed", "caught in our own audit; F1 0.23 to 0.98"),
        ("Compete on forecasting", "beats persistence at 9 of 10 horizons"),
        ("Second channel", "unsupervised surprise covers the gap"),
        ("Reported, not hidden", "detection transfers; progression does not"),
    ])

    add_band(slide, rx, BOTTOM - 1.02, rw, 0.44,
             "Stage forecast macro-F1 by horizon", fill=NAVY, size=10.5)
    add_table(slide, rx, BOTTOM - 0.56, rw, 0.56,
              [
                  ["Horizon", "+2", "+4", "+6", "+10"],
                  ["World model (ours)", "0.583", "0.642", "0.624", "0.524"],
                  ["“nothing changes” baseline", "0.473", "0.474",
                   "0.455", "0.436"],
              ],
              col_widths=[2.20, 0.96, 0.96, 0.96, 0.96], highlight_row=1, size=9)


def build_impact(slide) -> None:
    clear(find(slide, "TextBox 8"))

    add_band(slide, LEFT, TOP, WIDTH, 0.46,
             "Shifting the defender from reacting after compromise to acting "
             "during the kill chain", size=12.5)

    lw = 6.00
    rx = LEFT + lw + 0.40
    rw = RIGHT - rx

    section(slide, LEFT, TOP + 0.58, lw, "Potential impact on the target audience")
    add_hub(slide, LEFT + lw / 2, TOP + 3.00,
            ("AttackForecast", "10 minutes of warning"),
            [
                ("SOC analysts", "a ranked queue, not an alert flood"),
                ("Critical infrastructure", "power, banking, telecom, transport"),
                ("CERT-In", "national response teams"),
                ("Enterprises", "no new hardware to buy"),
                ("Auditors", "every alert is challengeable"),
            ], rx=2.10, ry=1.52)

    section(slide, rx, TOP + 0.58, rw, "Benefits of AttackForecast")
    add_ring(slide, rx + rw / 2, TOP + 3.00, [
        ("Early warning", "up to 10 minutes before the chain closes"),
        ("Low noise", "4.9% false alarms across 13 captures"),
        ("Auditable", "replay shows what the model knew, minute by minute"),
        ("Air-gap ready", "nothing leaves the network"),
        ("Zero licence cost", "open source, runs on existing hardware"),
        ("Playbook-ready", "ATT&CK stage names map onto what SOCs already use"),
    ])

    add_band(slide, LEFT, BOTTOM - 0.44, WIDTH, 0.44,
             "Every claim on this deck is reproducible from the repository: one "
             "command prepares the data, one trains, one benchmarks.",
             fill=GREEN, size=11.5)


def build_references(slide) -> None:
    clear(find(slide, "TextBox 8"))

    section(slide, LEFT, TOP, WIDTH, "Datasets, standards and prior work")

    cols = [
        ("CTU-13 Dataset",
         ["Stratosphere IPS, CTU Prague",
          "13 real botnet captures, CC BY",
          "Garcia et al., Computers &",
          "Security, 2014"]),
        ("MITRE ATT&CK",
         ["Enterprise Matrix",
          "TA0043 Reconnaissance",
          "TA0001 / TA0008 / TA0011",
          "TA0010 Exfiltration"]),
        ("World models",
         ["Ha & Schmidhuber,",
          "World Models, NeurIPS 2018",
          "Hafner et al., PlaNet, ICML 2019",
          "Hafner et al., Dreamer, ICLR 2020"]),
        ("Explainability",
         ["Sundararajan et al.,",
          "Integrated Gradients, ICML 2017",
          "Vaswani et al., Attention Is",
          "All You Need, NeurIPS 2017"]),
        ("National context",
         ["NCIIPC - nciipc.gov.in",
          "CIC-IDS2017 / CSE-CIC-IDS2018",
          "reviewed; access now gated",
          "behind registration"]),
    ]
    n = len(cols)
    gap = 0.16
    cw = (WIDTH - (n - 1) * gap) / n
    for i, (head, body) in enumerate(cols):
        colour = PALETTE[i % len(PALETTE)]
        cx = LEFT + i * (cw + gap)
        shp(slide, MSO_SHAPE.OVAL, cx + cw / 2 - 0.24, TOP + 0.50, 0.48, 0.48,
            fill=colour, lines=[str(i + 1)], size=13, bold=True)
        connect(slide, cx + cw / 2, TOP + 0.98, cx + cw / 2, TOP + 1.16,
                color=colour, width=1.5, dash=False)
        down_pentagon(slide, cx, TOP + 1.16, cw, 1.98, fill=colour)
        label(slide, cx + 0.10, TOP + 1.30, cw - 0.20, 1.50,
              [(head, 10, True)] + [(b, 7.5, False) for b in body],
              color=WHITE, spacing=2)

    section(slide, LEFT, TOP + 3.44, WIDTH, "Built with")
    add_notes(slide, LEFT, TOP + 3.92, WIDTH, [
        "Python 3.13", "PyTorch", "Scapy", "pandas + NumPy",
        "scikit-learn", "FastAPI", "Integrated Gradients",
    ], h=0.68)

    add_band(slide, LEFT, BOTTOM - 0.92, WIDTH, 0.44,
             "Problem Statement 26153 (NTRO)  -  AI based Network Attack "
             "Forecasting from Network Traffic Data  -  "
             "Blockchain & Cybersecurity",
             fill=NAVY, size=11)
    add_band(slide, LEFT, BOTTOM - 0.44, WIDTH, 0.44,
             "Source, data preparation and benchmark: "
             "github.com/HowSuyash/AttackForecast",
             fill=GREEN, size=11)


# --------------------------------------------------------------------------


def main() -> None:
    default_template = (Path.home() / "Downloads" /
                        "SIH2026-IDEA-Presentation-Format.pptx")
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--template", type=Path, default=default_template)
    parser.add_argument("--out", type=Path,
                        default=Path(__file__).parent.parent /
                        "SIH2026_PS26153_Presentation.pptx")
    parser.add_argument("--team-id", default="",
                        help="omitted from the title slide when empty")
    parser.add_argument("--team-name", default="Git with It")
    args = parser.parse_args()

    if not args.template.exists():
        raise SystemExit(f"Template not found: {args.template}")

    prs = Presentation(str(args.template))

    build_title(prs.slides[0], args.team_id, args.team_name)
    build_idea(prs.slides[1])
    build_technical(prs.slides[2])
    build_feasibility(prs.slides[3], Path(__file__).parent / "images")
    build_impact(prs.slides[4])
    build_references(prs.slides[5])

    # The template caps the deck at six slides including the title, and says
    # the instructions slide may be removed before upload.
    delete_slide(prs, 6)

    # The team-name ovals are template furniture; fill them so the deck does
    # not go out with placeholder text on every slide.
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name.startswith("Oval") and shape.has_text_frame:
                if "Your Team Name" in shape.text_frame.text:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = args.team_name
                            run.font.size = Pt(9)

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.out))
    print(f"wrote {args.out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
